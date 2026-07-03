"""
PPTX export — renders draft_proposal.json into a CYLAD-format deck.

    python export_pptx.py draft_proposal.json \
        --template "~/Downloads/202006 - PF - DAIS - REX Harmonie - Proposition CYLAD V3.pptx" \
        --out draft_proposal.pptx [--slide-library ~/Downloads] [--force-draft]

Design (template-first, never style from scratch):
  - Slides are built ONLY from the template's own layouts, keyed by name:
      cover      -> 'Front cover basic'      (title / subtitle / city-date)
      agenda     -> 'Content with 1 text box' (bullet list of chapters)
      content    -> 'Content with 1 text box' (action title, kicker idx15,
                    body idx16, comment strip idx11)
      back cover -> 'Back cover light'
  - Generated sections come in as SlideSpec lists: from GeminiSlidifier when
    the caller supplies one, else the markdown fallback here (headers ->
    kickers, list items -> bullets, paragraphs -> bullets).
  - retrieve_only sections (team, credentials): each entry's ORIGINAL slide is
    copied verbatim from the source deck when it can be found in the slide
    library (matched by file_name from credential_refs); otherwise a text
    one-pager with the entry's provenance link.
  - Export gate: exportable=False refuses unless force_draft, and then every
    slide carries a DRAFT watermark. Escalated sections are only included in
    watermarked decks, flagged in the kicker.
"""

from __future__ import annotations

import argparse
import copy
import json
import pathlib
import re
import sys
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt
from pydantic import BaseModel, Field

# ---------------------------------------------------------------------------
# Slide schema — what the slidifier (Gemini or fallback) produces
# ---------------------------------------------------------------------------


class Bullet(BaseModel):
    text: str
    level: int = Field(0, ge=0, le=2)


class SlideSpec(BaseModel):
    action_title: str
    kicker: str = ""            # subtitle strip above the body (idx 15)
    bullets: list[Bullet] = Field(default_factory=list)
    comment: str = ""           # footnote strip (idx 11)


_MAX_BULLETS = 9  # continuation slide past this


# ---------------------------------------------------------------------------
# Markdown fallback slidify — no LLM: headers structure, bullets carry over
# ---------------------------------------------------------------------------


_MD_INLINE = re.compile(r"\*\*(.+?)\*\*|\*(.+?)\*|\[(.+?)\]\([^)]*\)|`(.+?)`")


def _plain(md: str) -> str:
    """Strip inline markdown, keep the text."""
    return _MD_INLINE.sub(lambda m: next(g for g in m.groups() if g is not None),
                          md).strip()


def fallback_slidify(title: str, md: str) -> list[SlideSpec]:
    """Deterministic markdown -> slides. One slide per top heading run;
    list items become bullets (indent -> level), paragraphs become bullets."""
    slides: list[SlideSpec] = []
    cur = SlideSpec(action_title=title)
    for raw in md.splitlines():
        line = raw.rstrip()
        if not line.strip():
            continue
        h = re.match(r"^(#{2,4})\s+(.*)", line)
        li = re.match(r"^(\s*)[-*]\s+(.*)", line)
        if h:
            text = _plain(h.group(2))
            if len(h.group(1)) == 2 and not cur.bullets:
                cur.action_title = text          # section's own H2 = the title
            elif len(h.group(1)) == 2:
                slides.append(cur)               # new top heading = new slide
                cur = SlideSpec(action_title=text)
            else:
                cur.bullets.append(Bullet(text=text, level=0))
        elif li:
            level = min(len(li.group(1)) // 2, 2)
            cur.bullets.append(Bullet(text=_plain(li.group(2)), level=level))
        else:
            cur.bullets.append(Bullet(text=_plain(line), level=0))
    slides.append(cur)
    return [s for s in split_overlong(slides) if s.bullets or s.action_title]


def split_overlong(slides: list[SlideSpec]) -> list[SlideSpec]:
    """Continuation slides past _MAX_BULLETS — applied to every spec source
    (the LLM slidifier does not reliably respect its bullet budget)."""
    out: list[SlideSpec] = []
    for s in slides:
        chunks = [s.bullets[i:i + _MAX_BULLETS]
                  for i in range(0, len(s.bullets), _MAX_BULLETS)] or [[]]
        for j, chunk in enumerate(chunks):
            out.append(SlideSpec(
                action_title=s.action_title + (" (cont'd)" if j else ""),
                kicker=s.kicker if j == 0 else "",
                bullets=chunk, comment=s.comment if j == 0 else ""))
    return out


# ---------------------------------------------------------------------------
# Template helpers
# ---------------------------------------------------------------------------


def _layout(prs: Presentation, name: str):
    for lay in prs.slide_masters[0].slide_layouts:
        if lay.name == name:
            return lay
    raise KeyError(f"template has no layout named {name!r} — "
                   f"got {[l.name for l in prs.slide_masters[0].slide_layouts]}")


def _ph(slide, idx: int):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def _set_text(slide, idx: int, text: str, *, hyperlink: str | None = None) -> None:
    ph = _ph(slide, idx)
    if ph is None or not text:
        return
    tf = ph.text_frame
    tf.text = text
    if hyperlink:
        for run in tf.paragraphs[0].runs:
            run.hyperlink.address = hyperlink


def _fill_body(slide, idx: int, bullets: Iterable[Bullet]) -> None:
    ph = _ph(slide, idx)
    if ph is None:
        return
    tf = ph.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = b.text
        p.level = b.level


def _watermark(slide, prs: Presentation) -> None:
    box = slide.shapes.add_textbox(
        Emu(int(prs.slide_width * 0.1)), Emu(int(prs.slide_height * 0.42)),
        Emu(int(prs.slide_width * 0.8)), Emu(int(prs.slide_height * 0.16)))
    box.rotation = 345
    tf = box.text_frame
    tf.text = "DRAFT — PENDING HUMAN REVIEW"
    run = tf.paragraphs[0].runs[0]
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xC0, 0xC0, 0xC0)


# ---------------------------------------------------------------------------
# Verbatim slide copy — clone a slide (shapes + image rels) from a source deck
# ---------------------------------------------------------------------------


def copy_slide(src_path: pathlib.Path, slide_number: int, dst: Presentation):
    """Append source slide `slide_number` (1-based) to dst, cloning shapes and
    re-binding image relationships. Charts/OLE keep their drawn appearance only
    if they are pictures; live objects beyond images are not carried over."""
    src = Presentation(str(src_path))
    if not (1 <= slide_number <= len(src.slides)):
        raise IndexError(f"{src_path.name} has {len(src.slides)} slides, "
                         f"asked for {slide_number}")
    s = src.slides[slide_number - 1]

    blank = _layout(dst, "Titre seul")  # least chrome; copied shapes cover it
    new = dst.slides.add_slide(blank)
    for shape in list(new.shapes):      # drop the layout's empty placeholders
        shape._element.getparent().remove(shape._element)

    # map old image rIds -> new ones, then rewrite refs in the copied XML.
    # Foreign parts keep their partname on save, so a name already used in the
    # destination package (e.g. both decks have image11.emf) would produce a
    # corrupt zip with duplicate entries — rename to a free slot first.
    from pptx.opc.packuri import PackURI
    taken = {str(p.partname) for p in dst.part.package.iter_parts()}
    rid_map: dict[str, str] = {}
    for rel in s.part.rels.values():
        if "image" in rel.reltype and not rel.is_external:
            image_part = rel.target_part
            name = str(image_part.partname)
            if name in taken:
                stem, ext = name.rsplit(".", 1)
                n = 1
                while f"/ppt/media/pbcopy{n}.{ext}" in taken:
                    n += 1
                image_part.partname = PackURI(f"/ppt/media/pbcopy{n}.{ext}")
            taken.add(str(image_part.partname))
            new_rid = new.part.relate_to(image_part, rel.reltype)
            rid_map[rel.rId] = new_rid

    for shape in s.shapes:
        el = copy.deepcopy(shape._element)
        xml = el.xml if isinstance(el.xml, str) else el.xml.decode()
        if rid_map:
            for old, newr in rid_map.items():
                xml = re.sub(f'(r:(?:embed|link)=")({old})(")',
                             rf"\g<1>{newr}\g<3>", xml)
            from lxml import etree
            el = etree.fromstring(xml.encode())
        new.shapes._spTree.append(el)
    return new


def find_in_library(file_name: str, library: list[pathlib.Path]) -> pathlib.Path | None:
    for root in library:
        root = root.expanduser()
        if not root.is_dir():
            continue
        hit = next((f for f in root.rglob(file_name)), None)
        if hit:
            return hit
    return None


# ---------------------------------------------------------------------------
# Deck assembly
# ---------------------------------------------------------------------------


_ENTRY_HEAD = re.compile(r"^\*\*\[?(.+?)\]?(?:\((\S+)\))?\*\*$")


def _verbatim_entries(content: str) -> list[tuple[str, str | None, str]]:
    """Parse worker._verbatim_content back into (label, url, text) entries."""
    out = []
    for block in content.split("\n\n---\n\n"):
        lines = block.strip().splitlines()
        label, url, body = "", None, block.strip()
        if lines:
            m = _ENTRY_HEAD.match(lines[0].strip())
            if m:
                label = m.group(1).split("](")[0]
                url = m.group(2) or (m.group(1).split("](")[1].rstrip(")")
                                     if "](" in m.group(1) else None)
                body = "\n".join(lines[1:]).strip()
        out.append((label, url, body))
    return out


def export_pptx(
    payload: dict,
    *,
    template: pathlib.Path,
    out: pathlib.Path,
    slide_library: list[pathlib.Path] | None = None,
    slides_by_section: dict[str, list[SlideSpec]] | None = None,
    force_draft: bool = False,
) -> pathlib.Path:
    """Render the cli.py JSON payload to a deck. `slides_by_section` lets the
    caller pass Gemini-slidified specs; anything missing uses the fallback."""
    assembled = payload["assembled"]
    brief = payload.get("brief", {})
    exportable = assembled.get("exportable", False)
    if not exportable and not force_draft:
        raise RuntimeError(
            "review queue is not clear — refusing to export. "
            "Pass force_draft=True for a watermarked draft.")
    watermark = not exportable

    prs = Presentation(str(template))
    # keep the template's masters/layouts, drop its slides
    for i in range(len(prs.slides) - 1, -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]

    slides_by_section = slides_by_section or {}
    library = slide_library or []
    new_slides = []

    sections = assembled["sections"]
    escalated = {r["section"] for r in assembled.get("review_queue", [])}

    # cover
    cover_item = next((s for s in sections if s["section"] == "cover"), None)
    cover = prs.slides.add_slide(_layout(prs, "Front cover basic"))
    _set_text(cover, 0, (cover_item or {}).get("title")
              or f"{brief.get('client') or 'Client'} — Proposal")
    _set_text(cover, 1, " — ".join(filter(None, [brief.get("client"),
                                                 brief.get("service_line")])))
    import datetime
    _set_text(cover, 12, datetime.date.today().strftime("%d %B %Y"))
    new_slides.append(cover)

    # agenda
    chapters = [s["title"] for s in sections
                if s["section"] not in ("cover", "next_steps")]
    agenda = prs.slides.add_slide(_layout(prs, "Content with 1 text box"))
    _set_text(agenda, 0, "Agenda")
    _fill_body(agenda, 16, [Bullet(text=t) for t in chapters])
    new_slides.append(agenda)

    for s in sections:
        name, mode, status = s["section"], s.get("mode"), s["status"]
        if name == "cover":
            continue
        flagged = name in escalated or status.startswith("escalated")

        if mode == "retrieve_only" and not flagged:
            for label, url, body in _verbatim_entries(s.get("content", "")):
                src = find_in_library(label.split(" · ")[0], library) if label else None
                slide_no = re.search(r"slide (\d+)", label or "")
                if src and slide_no:
                    new_slides.append(
                        copy_slide(src, int(slide_no.group(1)), prs))
                    continue
                sl = prs.slides.add_slide(_layout(prs, "Content with 1 text box"))
                _set_text(sl, 0, s["title"])
                _set_text(sl, 15, label or "", hyperlink=url)
                _fill_body(sl, 16, [Bullet(text=t) for t in body.splitlines() if t.strip()])
                if url:
                    _set_text(sl, 11, url)
                new_slides.append(sl)
            continue

        specs = slides_by_section.get(name) or fallback_slidify(
            s["title"], s.get("content", ""))
        for spec in split_overlong(specs):
            sl = prs.slides.add_slide(_layout(prs, "Content with 1 text box"))
            _set_text(sl, 0, spec.action_title)
            kicker = spec.kicker
            if flagged:
                kicker = f"ESCALATED — HUMAN REVIEW REQUIRED ({status}). {kicker}"
            _set_text(sl, 15, kicker)
            _fill_body(sl, 16, spec.bullets)
            _set_text(sl, 11, spec.comment)
            new_slides.append(sl)

    back = prs.slides.add_slide(_layout(prs, "Back cover light"))
    _set_text(back, 12, "Thank you")
    _set_text(back, 11, "for your attention")
    new_slides.append(back)

    if watermark:
        for sl in new_slides:
            _watermark(sl, prs)

    out = out.expanduser()
    prs.save(str(out))
    return out


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------


def main() -> int:
    ap = argparse.ArgumentParser(description="Export draft_proposal.json to PPTX")
    ap.add_argument("payload", type=pathlib.Path, nargs="?",
                    default=pathlib.Path("draft_proposal.json"))
    ap.add_argument("--template", type=pathlib.Path, required=True)
    ap.add_argument("--out", type=pathlib.Path,
                    default=pathlib.Path("draft_proposal.pptx"))
    ap.add_argument("--slide-library", type=pathlib.Path, action="append",
                    default=[], help="Folder(s) with the source decks for "
                    "verbatim team/credential slides (searched recursively)")
    ap.add_argument("--force-draft", action="store_true",
                    help="Export despite a pending review queue (watermarked)")
    args = ap.parse_args()

    payload = json.loads(args.payload.read_text())
    out = export_pptx(payload, template=args.template.expanduser(), out=args.out,
                      slide_library=args.slide_library or None,
                      force_draft=args.force_draft)
    print(f"Wrote {out}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
